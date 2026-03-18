"""
PPTX Normalization.

Handles Microsoft PowerPoint presentations (.pptx format).

Processing rules:

  1. SLIDE ORDER
     Slides are processed in presentation order (slide index 1..N).
     slide_number is 1-based and stored on every TextBlock produced
     from that slide.

  2. SHAPE PROCESSING ORDER
     Within each slide, shapes are processed in z-order (shape index order).
     Title shapes are extracted first regardless of their z-order:
       - A shape is a title if placeholder_format is not None
         and placeholder_format.idx == 0
     All remaining shapes are processed in z-order after the title.

  3. TEXT EXTRACTION
     For each shape with has_text_frame=True:
       - Iterate shape.text_frame.paragraphs in order.
       - Join all run text within each paragraph to form paragraph text.
       - Skip paragraphs whose text is whitespace-only.
     Title shapes produce a TextBlock with heading_level=1.
     All other text shapes produce TextBlocks with heading_level=None.

  4. SPEAKER NOTES
     Each slide's notes_slide.notes_text_frame is extracted.
     Notes produce a TextBlock with is_speaker_notes=True.
     Empty notes slides are skipped.

  5. EMBEDDED TABLES
     Shapes with has_table=True are extracted as ExtractedTable objects.
     First row is the header row. All cell values cast to str and stripped.
     Empty tables (no rows) are skipped.

  6. IMAGE-ONLY SLIDES
     A slide that produces zero text content (no text from any shape,
     no notes, no table) gets a single placeholder TextBlock:
       is_placeholder=True
       text = "[Image-only slide — no text content extracted]"
     This signals to downstream processing that OCR may be needed.

Output: NormalizedDocument with TextBlocks and ExtractedTables.
"""
from __future__ import annotations

import logging
from typing import Optional
from uuid import UUID, uuid4

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from packages.core.schemas.document import (
    DocumentMetadata,
    ExtractedTable,
    NormalizedDocument,
    TextBlock,
)

logger = logging.getLogger(__name__)


def normalize_pptx(
    content: bytes,
    document_id: UUID,
    engagement_id: UUID,
    filename: str,
    vdr_folder_path: Optional[str] = None,
) -> NormalizedDocument:
    """
    Normalize a PPTX file to a NormalizedDocument.

    Args:
        content:          Raw PPTX bytes.
        document_id:      UUID of the document record.
        engagement_id:    UUID of the engagement.
        filename:         Original filename (used in citations).
        vdr_folder_path:  Optional VDR folder path for metadata.

    Returns:
        NormalizedDocument with all TextBlocks and ExtractedTables.
    """
    import io
    prs = Presentation(io.BytesIO(content))

    text_blocks: list[TextBlock] = []
    tables: list[ExtractedTable] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1  # 1-based

        slide_blocks, slide_tables = _process_slide(
            slide, slide_number, document_id
        )
        text_blocks.extend(slide_blocks)
        tables.extend(slide_tables)

    return NormalizedDocument(
        document_id=document_id,
        engagement_id=engagement_id,
        text_blocks=text_blocks,
        tables=tables,
        metadata=DocumentMetadata(
            filename=filename,
            file_type="pptx",
            page_count=len(prs.slides),
            vdr_folder_path=vdr_folder_path,
            section_path=[],
        ),
    )


# ─── Slide processing ─────────────────────────────────────────────────────────

def _process_slide(
    slide,
    slide_number: int,
    document_id: UUID,
) -> tuple[list[TextBlock], list[ExtractedTable]]:
    """
    Process one slide into TextBlocks and ExtractedTables.

    Returns (text_blocks, tables).
    If the slide has no extractable text content, a placeholder TextBlock
    is added for the image-only slide.
    """
    text_blocks: list[TextBlock] = []
    tables: list[ExtractedTable] = []

    # Separate title shapes from other shapes
    title_shapes = []
    other_shapes = []
    for shape in slide.shapes:
        if _is_title_shape(shape):
            title_shapes.append(shape)
        else:
            other_shapes.append(shape)

    # Process title shapes first
    for shape in title_shapes:
        blocks = _extract_text_blocks(shape, slide_number, document_id, is_title=True)
        text_blocks.extend(blocks)

    # Process remaining shapes in z-order
    for shape in other_shapes:
        if shape.has_text_frame:
            blocks = _extract_text_blocks(shape, slide_number, document_id, is_title=False)
            text_blocks.extend(blocks)

        if shape.has_table:
            table = _extract_table(shape.table, slide_number, document_id)
            if table is not None:
                tables.append(table)

    # Extract speaker notes
    notes_block = _extract_speaker_notes(slide, slide_number, document_id)
    if notes_block is not None:
        text_blocks.append(notes_block)

    # Image-only slide detection:
    # If no text came from shapes (only consider non-notes blocks)
    non_notes = [b for b in text_blocks if not b.is_speaker_notes]
    if not non_notes and not tables:
        text_blocks.append(
            TextBlock(
                block_id=uuid4(),
                document_id=document_id,
                slide_number=slide_number,
                page_number=None,
                heading_level=None,
                text="[Image-only slide — no text content extracted]",
                is_placeholder=True,
            )
        )
        logger.info("Slide %d is image-only (placeholder inserted)", slide_number)

    return text_blocks, tables


# ─── Shape helpers ────────────────────────────────────────────────────────────

def _is_title_shape(shape) -> bool:
    """
    Return True if the shape is a title placeholder.
    Title placeholders have placeholder_format.idx == 0.
    """
    try:
        ph = shape.placeholder_format
        return ph is not None and ph.idx == 0
    except Exception:
        return False


def _extract_text_blocks(
    shape,
    slide_number: int,
    document_id: UUID,
    is_title: bool,
) -> list[TextBlock]:
    """
    Extract TextBlocks from a shape's text frame.

    Title shapes produce heading_level=1 blocks.
    Other shapes produce heading_level=None blocks.
    Empty paragraphs are skipped.
    """
    if not shape.has_text_frame:
        return []

    blocks: list[TextBlock] = []
    heading_level = 1 if is_title else None

    for para in shape.text_frame.paragraphs:
        # Join all run text in this paragraph
        text = "".join(run.text for run in para.runs if run.text)
        if not text.strip():
            continue
        blocks.append(
            TextBlock(
                block_id=uuid4(),
                document_id=document_id,
                slide_number=slide_number,
                page_number=None,
                heading_level=heading_level,
                text=text.strip(),
            )
        )

    return blocks


def _extract_speaker_notes(
    slide,
    slide_number: int,
    document_id: UUID,
) -> Optional[TextBlock]:
    """
    Extract speaker notes from a slide's notes_slide.

    Returns a TextBlock with is_speaker_notes=True, or None if notes are empty.
    """
    try:
        notes_slide = slide.notes_slide
        if notes_slide is None:
            return None
        notes_tf = notes_slide.notes_text_frame
        if notes_tf is None:
            return None
        notes_text = notes_tf.text.strip()
        if not notes_text:
            return None
        return TextBlock(
            block_id=uuid4(),
            document_id=document_id,
            slide_number=slide_number,
            page_number=None,
            heading_level=None,
            text=notes_text,
            is_speaker_notes=True,
        )
    except Exception as exc:
        logger.debug("Could not extract notes from slide %d: %s", slide_number, exc)
        return None


def _extract_table(
    table,
    slide_number: int,
    document_id: UUID,
) -> Optional[ExtractedTable]:
    """
    Extract an ExtractedTable from a PowerPoint table shape.

    First row is the header row.
    All cell values cast to str and stripped.
    Returns None for empty tables.
    """
    if not table.rows:
        return None

    raw_rows: list[list[str]] = []
    for row in table.rows:
        raw_rows.append([cell.text.strip() for cell in row.cells])

    if not raw_rows:
        return None

    expected_cols = len(raw_rows[0])
    if expected_cols == 0:
        return None

    # Enforce column consistency
    clean_rows = [r for r in raw_rows if len(r) == expected_cols]
    if not clean_rows:
        return None

    headers = clean_rows[0]
    if not any(h.strip() for h in headers):
        return None

    return ExtractedTable(
        table_id=uuid4(),
        document_id=document_id,
        page_number=None,
        headers=headers,
        rows=clean_rows[1:],
        table_type="generic",
        is_sampled=False,
        is_malformed=False,
    )
